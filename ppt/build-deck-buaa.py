#!/usr/bin/env python3
"""Build a 10-page CleanTest-Agent deck on top of the BUAA template.

Source of truth for content: ppt/slides-marp.md (which is itself
content-equivalent to ppt/slides.md / report/main.tex).

Output: ppt/CleanTest-Agent-BUAA.pptx

The script clones the school-provided 北航PPT版式1.pptx, replaces
its single placeholder slide with 10 information-dense pages, all
drawn from the master's existing slideLayouts so the BUAA branding
(header band, page footer, fonts, colours) is preserved verbatim.
"""

from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
TEMPLATE = HERE / "北航PPT版式1.pptx"
OUTPUT = HERE / "CleanTest-Agent-BUAA.pptx"

# ----- Colour palette (matches the Marp deck, derived from BUAA blue) -----
BUAA_BLUE = RGBColor(0x1A, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)
ACCENT_RED = RGBColor(0xB9, 0x1C, 0x1C)
NEUTRAL_DARK = RGBColor(0x11, 0x11, 0x11)
NEUTRAL_GREY = RGBColor(0x55, 0x55, 0x55)
BG_LIGHT = RGBColor(0xEF, 0xF6, 0xFF)
BG_BORDER = RGBColor(0x25, 0x63, 0xEB)


# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------

def remove_existing_slides(prs: Presentation) -> None:
    """Strip the placeholder slide that ships with the BUAA template."""
    sldIdLst = prs.slides._sldIdLst  # noqa: SLF001 (private API, stable)
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def add_slide(prs: Presentation, layout_idx: int):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def set_text(tf, text: str, *, size: int = 18, bold: bool = False,
             colour: RGBColor = NEUTRAL_DARK, align=None) -> None:
    """Replace the entire text frame content with a single line."""
    tf.clear()
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour


def add_paragraph(tf, text: str, *, size: int = 16, bold: bool = False,
                  colour: RGBColor = NEUTRAL_DARK, level: int = 0,
                  bullet: bool = True, space_before: int = 2) -> None:
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(space_before)
    if bullet:
        # The BUAA layout's body placeholder already provides bullets; we
        # simply add level-0 text here. To suppress bullets explicitly,
        # caller can pass bullet=False (handled via a no-op for now since
        # the underlying layout decides).
        pass
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour


def set_title(slide, title: str, *, size: int = 28) -> None:
    ph = slide.shapes.title
    if ph is None:
        return
    set_text(ph.text_frame, title, size=size, bold=True, colour=BUAA_BLUE)


def fill_body(ph, lines: list[tuple[str, dict]] | list[str]) -> None:
    """Fill a body placeholder. `lines` is a list of either
    `str` (default formatting) or `(text, kwargs)` tuples."""
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for item in lines:
        text, kwargs = (item, {}) if isinstance(item, str) else item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = kwargs.get("level", 0)
        run = p.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = Pt(kwargs.get("size", 14))
        run.font.bold = kwargs.get("bold", False)
        run.font.color.rgb = kwargs.get("colour", NEUTRAL_DARK)


# Add a coloured highlight box (manual rectangle outside the placeholders).
def add_highlight_box(slide, *, left: float, top: float, width: float,
                      height: float, lines: list[tuple[str, dict]]) -> None:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = BG_LIGHT
    box.line.color.rgb = BG_BORDER
    box.line.width = Pt(1.0)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    first = True
    for text, kwargs in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = kwargs.get("align", PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = Pt(kwargs.get("size", 13))
        run.font.bold = kwargs.get("bold", False)
        run.font.color.rgb = kwargs.get("colour", NEUTRAL_DARK)


def add_table(slide, *, left: float, top: float, width: float, height: float,
              data: list[list[str]],
              col_widths: list[float] | None = None,
              header_size: int = 13, body_size: int = 12,
              highlight_rows: list[int] | None = None,
              highlight_cols: list[int] | None = None) -> None:
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(
        rows, cols,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tbl = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.name = "Calibri"
            run.font.size = Pt(header_size if r == 0 else body_size)
            run.font.bold = (r == 0)
            if r == 0:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = BUAA_BLUE
            else:
                if highlight_rows and r in highlight_rows:
                    run.font.bold = True
                    run.font.color.rgb = ACCENT_RED
                if highlight_cols and c in highlight_cols and r != 0:
                    run.font.bold = True
                    run.font.color.rgb = ACCENT_RED


# --------------------------------------------------------------------------
# Build the deck
# --------------------------------------------------------------------------

def build() -> None:
    prs = Presentation(str(TEMPLATE))
    remove_existing_slides(prs)

    # ====== Slide 1 — Cover (layout 0: 标题幻灯片) ============================
    s = add_slide(prs, 0)
    # The cover layout has NO text placeholders (it is purely decorative,
    # carrying the BUAA logo at the top-left and the campus silhouette at
    # the bottom). We drop a manual title + subtitle text box centred in
    # the empty space between them.
    title_tb = s.shapes.add_textbox(Inches(0.92), Inches(2.20), Inches(11.5), Inches(1.20))
    set_text(title_tb.text_frame, "CleanTest-Agent",
             size=54, bold=True, colour=BUAA_BLUE)
    sub = s.shapes.add_textbox(Inches(0.92), Inches(3.50), Inches(11.5), Inches(2.5))
    tf = sub.text_frame
    tf.word_wrap = True
    add_paragraph(tf,
        "A Multi-Agent Skill-Orchestrated System for Unit Test Training Data Quality Assurance",
        size=22, bold=True, colour=ACCENT_BLUE)
    add_paragraph(tf, "", size=8)
    add_paragraph(tf, "Yong Yang  ·  School of Software, Beihang University",
                  size=18, colour=NEUTRAL_DARK, space_before=10)
    add_paragraph(tf, "yang_qhd@buaa.edu.cn",
                  size=18, colour=ACCENT_BLUE)
    add_paragraph(tf, "Software Requirements Analysis & System Design  ·  Final Project",
                  size=16, colour=NEUTRAL_GREY, space_before=10)

    # ====== Slide 2 — Problem & Motivation (layout 3: 两栏内容) ===============
    s = add_slide(prs, 3)
    set_title(s, "Problem & Motivation: 43.52% of Methods2Test Is Noise")
    left_ph, right_ph = s.placeholders[1], s.placeholders[2]
    fill_body(left_ph, [
        ("Methods2Test: 780,944 test cases from 91,385 open-source Java projects.",
         {"size": 16}),
        ("43.52% contain noise (Zhang et al., FSE 2025 Distinguished Paper).",
         {"size": 16, "level": 0}),
        ("Filtering this noise improves downstream branch coverage by an average of 67% on Defects4J.",
         {"size": 16, "level": 0}),
        ("", {"size": 8}),
        ("Problems with the original CleanTest:",
         {"size": 15, "bold": True, "colour": ACCENT_RED}),
        ("- monolithic Python scripts; no test suite;",
         {"size": 14, "level": 1}),
        ("- not invokable from AI coding assistants;",
         {"size": 14, "level": 1}),
        ("- naive dictionary matching: ~30 min on 600 K samples.",
         {"size": 14, "level": 1}),
    ])
    # Right column: noise distribution as a tight in-cell table.
    fill_body(right_ph, [("Noise distribution (8 types):",
                          {"size": 14, "bold": True, "colour": BUAA_BLUE})])
    add_table(s,
        left=6.85, top=2.30, width=5.45, height=4.20,
        data=[
            ["Noise type", "% of corpus"],
            ["Unnecessary annotations", "41.64%"],
            ["No relevance", "12.70%"],
            ["Low coverage", "3.95%"],
            ["Ambiguous data type", "3.08%"],
            ["Syntax errors", "1.11%"],
            ["Empty exception handling", "0.68%"],
            ["Missing implementation", "0.34%"],
            ["Non-English literals", "0.16%"],
        ],
        col_widths=[3.40, 2.05],
        header_size=12, body_size=12,
        highlight_rows=[1],
    )

    # ====== Slide 3 — System Architecture (layout 1: 标题和内容) ==============
    s = add_slide(prs, 1)
    set_title(s, "System Architecture: Pipeline-of-Skills")
    body = s.placeholders[1]
    body.text_frame.clear()
    # ASCII-art pipeline as a monospaced text box, then bullets below.
    diag = s.shapes.add_textbox(Inches(0.6), Inches(1.10), Inches(12.1), Inches(2.85))
    diag_tf = diag.text_frame
    diag_tf.word_wrap = False
    diag_tf.margin_left = Inches(0.2)
    diag_tf.margin_top = Inches(0.05)
    diag_text = (
        "       User / Coding Assistant   (natural language)\n"
        "                       |\n"
        "       Orchestrator Skill   (cleantest-pipeline)\n"
        "         |               |               |\n"
        "         v               v               v\n"
        "     Filter 1         Filter 2         Filter 3\n"
        "     Syntax           Relevance        Coverage\n"
        "     AST + Aho-       Name match       Label scan (default) /\n"
        "     Corasick         + LLM            Qwen2.5-Coder-0.5B\n"
        "                      fallback         (model mode)\n"
        "         |               |               |\n"
        "         +---------------+---------------+\n"
        "                         v\n"
        "             Clean Dataset + Noise Report"
    )
    diag_tf.clear()
    p = diag_tf.paragraphs[0]
    run = p.add_run()
    run.text = diag_text
    run.font.name = "Menlo"
    run.font.size = Pt(11)
    run.font.color.rgb = NEUTRAL_DARK
    # Bullet summary below the diagram.
    bullets = s.shapes.add_textbox(Inches(0.6), Inches(4.30), Inches(12.1), Inches(2.50))
    btf = bullets.text_frame
    btf.word_wrap = True
    add_paragraph(btf,
        "4 composable Agent Skills sharing the SKILL.md protocol -- triggered by natural language from CodeBuddy, Claude Code, or Cursor.",
        size=14, colour=NEUTRAL_DARK)
    add_paragraph(btf,
        "Filter 3 ships TWO modes: deterministic JaCoCo label-scan (default) OR Qwen2.5-Coder-0.5B regression model when labels are missing.",
        size=14, colour=NEUTRAL_DARK)
    add_paragraph(btf,
        "Each filter is independently testable; 36-case pytest suite with Python 3.10 / 3.11 / 3.12 matrix.",
        size=14, colour=NEUTRAL_DARK)

    # ====== Slide 4 — Model-Driven Approach (layout 1) =======================
    s = add_slide(prs, 1)
    set_title(s, "Model-Driven Approach: Right Tool for the Right Subtask")
    s.placeholders[1].text_frame.clear()
    add_table(s,
        left=0.6, top=1.20, width=12.1, height=2.6,
        data=[
            ["", "Rules (deterministic)", "ML model (learned)", "LLM (semantic)"],
            ["What", "AST + Aho-Corasick over a 21,954-pattern dictionary",
             "Filter 3: JaCoCo label scan (default) or Qwen2.5-Coder-0.5B fine-tuned on 469 K rows",
             "DeepSeek-V4-Flash"],
            ["Coverage", "87.3% of samples decided here",
             "O(N) row scan in label mode",
             "Borderline ~12.7% only"],
            ["Time / sample", "~0.2 ms", "~0 (label mode)", "~3 s (API)"],
            ["$$ Cost", "0", "0 (label mode)", "~$4.5 / 593 K samples"],
        ],
        col_widths=[1.55, 3.85, 4.10, 2.60],
        header_size=12, body_size=11,
    )
    add_highlight_box(s,
        left=0.6, top=4.05, width=12.1, height=0.75,
        lines=[
            ("vs. pure LLM: rules handle 87.3% of samples free and instantly; "
             "the LLM is invoked only on the 12.7% it adds value to -- "
             "and never on the 21,954-pattern annotation dictionary it cannot memorise.",
             {"size": 12, "colour": BUAA_BLUE}),
        ])
    # Design patterns line.
    bottom = s.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.7))
    btf = bottom.text_frame
    btf.word_wrap = True
    add_paragraph(btf, "Five design patterns applied:",
                  size=14, bold=True, colour=ACCENT_RED)
    add_paragraph(btf, "  Strategy (filter selection)  |  Pipeline (sequential stages)  |  Facade (LLM-client abstraction)  |",
                  size=13, colour=NEUTRAL_DARK)
    add_paragraph(btf, "  Observer (NoiseReport accumulator)  |  Reflection (Filter 2 self-correction).",
                  size=13, colour=NEUTRAL_DARK)

    # ====== Slide 5 — Two Methodological Highlights (layout 4: 比较) =========
    s = add_slide(prs, 4)
    set_title(s, "Two Methodological Highlights over the Original CleanTest")
    # Left subtitle + body
    set_text(s.placeholders[1].text_frame,
             "Highlight 1 -- Reflection (Filter 2, opt-in)",
             size=15, bold=True, colour=ACCENT_RED)
    fill_body(s.placeholders[2], [
        ("Inspired by Lab 1's Reflection Agent pattern.", {"size": 13}),
        ("On IRRELEVANT verdict, apply a 5-rule checklist:", {"size": 13}),
        ("R1 Call Graph (Tufano 2022, extended)", {"size": 12, "level": 1}),
        ("R2 State Verification (Meszaros 2007)", {"size": 12, "level": 1}),
        ("R3 Behavior Verification (Meszaros 2007)", {"size": 12, "level": 1}),
        ("R4 Naming Equivalence", {"size": 12, "level": 1}),
        ("R5 Counterfactual", {"size": 12, "level": 1}),
        ("On 45 zero-overlap samples: 7/45 verdicts changed (15.6%); 5 rescued from false removal (11.1%); 6/7 verified correct (85.7% accuracy).",
         {"size": 12, "colour": BUAA_BLUE}),
        ("Default OFF -- additive on top of headline F1 = 0.965.",
         {"size": 12, "bold": True, "colour": ACCENT_RED}),
    ])
    # Right subtitle + body (Filter 3 model mode)
    set_text(s.placeholders[3].text_frame,
             "Highlight 2 -- Filter 3 Model Mode (replaces CodeGPT)",
             size=15, bold=True, colour=ACCENT_RED)
    s.placeholders[4].text_frame.clear()
    add_table(s,
        left=6.75, top=2.74, width=5.67, height=3.40,
        data=[
            ["Metric (held-out 46,921)", "This work (Qwen 0.5B)", "CodeGPT (Zhang 2025)"],
            ["MAE", "0.0309", "0.0798 (~2.6x)"],
            ["MSE", "0.0039", "0.0105 (~2.7x)"],
            ["RMSE", "0.0628", "--"],
            ["R-squared", "0.604", "--"],
            ["Pearson r / Spearman rho", "0.778 / 0.848", "--"],
            ["F1 @ tau = 0.10", "0.857", "--"],
            ["F1 @ tau = 0.15", "0.912", "--"],
        ],
        col_widths=[2.40, 1.80, 1.47],
        header_size=10, body_size=10,
        highlight_cols=[1],
    )
    # Caption under the right table.
    cap = s.shapes.add_textbox(Inches(6.75), Inches(6.20), Inches(5.67), Inches(0.55))
    add_paragraph(cap.text_frame,
        "Filter 3 now runs label-free on a single A800 80 GB (~3.32 h). The original CleanTest's hardest dependency on JaCoCo coverage labels is removed.",
        size=10, colour=NEUTRAL_GREY)

    # ====== Slide 6 — Aho-Corasick & Experiment Design (layout 4) ============
    s = add_slide(prs, 4)
    set_title(s, "Aho-Corasick Optimisation  +  Experiment Design")
    # Left: AC speedup
    set_text(s.placeholders[1].text_frame,
             "~11.5x pipeline speedup (filtering result unchanged)",
             size=14, bold=True, colour=ACCENT_RED)
    s.placeholders[2].text_frame.clear()
    add_table(s,
        left=0.92, top=2.74, width=5.64, height=2.70,
        data=[
            ["", "Naive", "Aho-Corasick"],
            ["Complexity", "O(N · K · L)", "O(N · (L + Z))"],
            ["K = 21,954 patterns", "scan all", "single pass"],
            ["Filter 1 (593,953)", "~30 min", "~1.6 min (~18.8x)"],
            ["Whole pipeline", "~31 min", "~2.6 min (~11.5x)"],
            ["Filtering result", "53.97% removed", "53.97% removed"],
        ],
        col_widths=[2.40, 1.70, 1.54],
        header_size=11, body_size=10,
        highlight_cols=[2],
    )
    add_highlight_box(s,
        left=0.92, top=5.55, width=5.64, height=1.10,
        lines=[
            ("Algorithmic optimisation at exactly the bottleneck the LLM cannot fix. "
             "Pure performance -- correctness unchanged.",
             {"size": 11, "colour": BUAA_BLUE}),
        ])
    # Right: 4 RQ + experiment setup
    set_text(s.placeholders[3].text_frame,
             "4 Research Questions  ·  real API + real GPU",
             size=14, bold=True, colour=ACCENT_RED)
    fill_body(s.placeholders[4], [
        ("RQ1. Reproduce CleanTest detection?  -> Yes (43.68% vs. 41.64%).",
         {"size": 12}),
        ("RQ2. Model-driven vs. pure LLM?  -> Real DeepSeek-V4-Flash exp.",
         {"size": 12}),
        ("RQ3. Aho-Corasick speedup?  -> ~11.5x (whole pipeline).",
         {"size": 12}),
        ("RQ4. Cost & speed?  -> Rules ~13,000-15,000x faster.",
         {"size": 12}),
        ("", {"size": 6}),
        ("Setup: DeepSeek-V4-Flash via Tencent Cloud TokenHub; 500 stratified samples (231 noise / 269 clean); 1,000 real API calls; T = 0.0.",
         {"size": 11, "colour": NEUTRAL_GREY}),
        ("Filter 3 model mode: A800 80 GB, bf16, batch 64, max_seq 512, lr 3e-5, 2 epochs (~3.32 h); 46,921-sample held-out test split.",
         {"size": 11, "colour": NEUTRAL_GREY}),
    ])

    # ====== Slide 7 — Results & Why Pure LLM Falls Short (layout 4) ==========
    s = add_slide(prs, 4)
    set_title(s, "Results: Rules + Selective LLM  >>  Pure LLM")
    # Left: results table
    set_text(s.placeholders[1].text_frame,
             "Real-API experiment (500 stratified Methods2Test samples)",
             size=14, bold=True, colour=ACCENT_RED)
    s.placeholders[2].text_frame.clear()
    add_table(s,
        left=0.92, top=2.74, width=5.64, height=2.40,
        data=[
            ["Method", "Prec.", "Recall", "F1", "Time"],
            ["Rule-based", "1.000", "1.000", "1.000", "0.11 s"],
            ["LLM zero-shot", "0.505", "0.221", "0.307", "1,487 s"],
            ["LLM few-shot", "0.534", "0.303", "0.387", "1,642 s"],
            ["Hybrid (ours)", "0.974", "0.956", "0.965", "< 60 s"],
        ],
        col_widths=[1.85, 0.95, 0.95, 0.95, 0.94],
        header_size=12, body_size=11,
        highlight_rows=[4],
    )
    add_highlight_box(s,
        left=0.92, top=5.30, width=5.64, height=1.40,
        lines=[
            ("F1 0.965 vs. 0.387  (+149%);  0.11 s vs. 1,642 s on 500 samples (~13,000-15,000x faster).",
             {"size": 11, "bold": True, "colour": ACCENT_RED}),
            ("LLM zero-shot misses 77.9% of noise; 87.8% of misses are unnecessary annotations.",
             {"size": 11, "colour": BUAA_BLUE}),
        ])
    # Right: case study
    set_text(s.placeholders[3].text_frame,
             "Case Study  --  Why pure LLM misjudges",
             size=14, bold=True, colour=ACCENT_RED)
    s.placeholders[4].text_frame.clear()
    code_box = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.75), Inches(2.74), Inches(5.67), Inches(1.45),
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    code_box.line.fill.background()
    ctf = code_box.text_frame
    ctf.margin_left = Inches(0.10)
    ctf.margin_top = Inches(0.05)
    ctf.clear()
    cp = ctf.paragraphs[0]
    cr = cp.add_run()
    cr.text = ("@PostMapping(path = \"/account/new\")\n"
               "public ResponseEntity<?> createAccount(\n"
               "    @Valid @RequestBody final AccountDto account, ...)\n"
               "{ ... }")
    cr.font.name = "Menlo"
    cr.font.size = Pt(11)
    cr.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

    add_table(s,
        left=6.75, top=4.30, width=5.67, height=1.40,
        data=[
            ["", "Rule-based", "LLM"],
            ["Verdict", "NOISE", "CLEAN (wrong)"],
            ["Reason", "AC matches @PostMapping", "looks like Spring controller"],
            ["Time", "< 0.01 ms", "~3,000 ms"],
        ],
        col_widths=[1.20, 2.00, 2.47],
        header_size=11, body_size=10,
        highlight_cols=[1],
    )
    cap = s.shapes.add_textbox(Inches(6.75), Inches(5.80), Inches(5.67), Inches(1.00))
    cap_tf = cap.text_frame
    cap_tf.word_wrap = True
    add_paragraph(cap_tf,
        "Root cause: LLM judges code quality, not training-data suitability.",
        size=10, colour=NEUTRAL_GREY)
    add_paragraph(cap_tf,
        "The 21,954-pattern dictionary cannot fit a single prompt window.",
        size=10, colour=NEUTRAL_GREY)

    # ====== Slide 8 — Software Engineering & CI/CD (layout 3: 两栏内容) =======
    s = add_slide(prs, 3)
    set_title(s, "Software Engineering Practices  +  CI / CD")
    fill_body(s.placeholders[1], [
        ("Requirements & Design",
         {"size": 16, "bold": True, "colour": BUAA_BLUE}),
        ("3 use cases  ·  9 functional + 6 non-functional requirements  ·  full traceability matrix.",
         {"size": 13}),
        ("UML use-case / activity / sequence diagrams.",
         {"size": 13}),
        ("Component diagram + class diagram.",
         {"size": 13}),
        ("5 design patterns:  Strategy / Pipeline / Facade / Observer / Reflection.",
         {"size": 13}),
        ("", {"size": 8}),
        ("Course-lab method integration",
         {"size": 16, "bold": True, "colour": BUAA_BLUE}),
        ("Lab 1 Reflection  ->  Filter 2 self-correction.",
         {"size": 13}),
        ("Lab 2 DSL validation  ->  rule-first, LLM only on residual.",
         {"size": 13}),
        ("Lab 3 formal verification  ->  AST structural checking.",
         {"size": 13}),
    ])
    fill_body(s.placeholders[2], [
        ("Testing & CI",
         {"size": 16, "bold": True, "colour": BUAA_BLUE}),
        ("36 pytest cases, 100% passing.",
         {"size": 13}),
        ("pytest + flake8 + mypy quality gate.",
         {"size": 13}),
        ("GitHub Actions matrix on Python 3.10 / 3.11 / 3.12.",
         {"size": 13}),
        ("", {"size": 8}),
        ("CD (publish.yml, tag-driven)",
         {"size": 16, "bold": True, "colour": BUAA_BLUE}),
        ("build  ->  twine check --strict",
         {"size": 13}),
        ("PyPI publish via OIDC Trusted Publisher (no API token)",
         {"size": 13}),
        ("sigstore keyless signing",
         {"size": 13}),
        ("attach signed wheel + sdist to GitHub Release",
         {"size": 13}),
        ("", {"size": 6}),
        ("Result:  pip install cleantest-agent  (PyPI v0.1.1).",
         {"size": 14, "bold": True, "colour": ACCENT_RED}),
    ])

    # ====== Slide 9 — Conclusion (layout 1) ==================================
    s = add_slide(prs, 1)
    set_title(s, "Conclusion -- Four Contributions")
    fill_body(s.placeholders[1], [
        ("1.  Skill-based architecture  --  4 composable Agent Skills via the SKILL.md protocol; natural-language trigger from CodeBuddy / Claude Code / Cursor.",
         {"size": 15}),
        ("2.  Hybrid rule + selective-LLM detection  --  F1 = 0.965 vs. pure-LLM 0.387 on real DeepSeek-V4-Flash experiments; ~13,000-15,000x faster.",
         {"size": 15}),
        ("3.  Aho-Corasick optimisation  --  ~11.5x pipeline speedup (~18.8x on Filter 1 alone), filtering result identical.",
         {"size": 15}),
        ("4.  Two methodological enhancements over the original CleanTest:",
         {"size": 15}),
        ("(a) Reflection (Filter 2, opt-in) -- 5-rule self-check rescues 5/45 (11.1%) borderline samples with 85.7% accuracy.",
         {"size": 14, "level": 1, "colour": NEUTRAL_DARK}),
        ("(b) Filter 3 model mode -- fine-tuned Qwen2.5-Coder-0.5B replaces CodeGPT with held-out MAE 0.0309 (~2.6x lower).",
         {"size": 14, "level": 1, "colour": NEUTRAL_DARK}),
    ])
    add_highlight_box(s,
        left=0.6, top=5.45, width=12.1, height=1.20,
        lines=[
            ("Systematic software design beats \"just ask the LLM\" -- supported by real-API experiments and a reproducible CI / CD pipeline.",
             {"size": 14, "bold": True, "colour": BUAA_BLUE, "align": PP_ALIGN.CENTER}),
            ("Released as v0.1.1 on PyPI:  pip install cleantest-agent  ·  sigstore-signed wheel attached to the v0.1.1 GitHub Release.",
             {"size": 12, "colour": NEUTRAL_GREY, "align": PP_ALIGN.CENTER}),
        ])

    # ====== Slide 10 — Thank You / Q&A (layout 0) ============================
    s = add_slide(prs, 0)
    title_tb = s.shapes.add_textbox(Inches(0.92), Inches(2.20), Inches(11.5), Inches(1.40))
    set_text(title_tb.text_frame, "Thank You",
             size=64, bold=True, colour=BUAA_BLUE)
    sub = s.shapes.add_textbox(Inches(0.92), Inches(3.70), Inches(11.5), Inches(2.50))
    tf = sub.text_frame
    tf.word_wrap = True
    add_paragraph(tf, "Questions?",
                  size=28, bold=True, colour=ACCENT_BLUE)
    add_paragraph(tf, "", size=8)
    add_paragraph(tf, "GitHub  ·  https://github.com/jimmy0717/cleantest-agent",
                  size=18, colour=NEUTRAL_DARK)
    add_paragraph(tf, "PyPI    ·  https://pypi.org/project/cleantest-agent/",
                  size=18, colour=NEUTRAL_DARK)
    add_paragraph(tf, "Contact ·  yang_qhd@buaa.edu.cn",
                  size=18, colour=NEUTRAL_DARK)

    # ----- Save -----
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}  ({OUTPUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
